"""Attachment helpers: save uploaded files and extract text from many document types."""

import base64
import io
import os
import re
import uuid
from pathlib import Path

import fitz  # PyMuPDF
from fastapi import HTTPException, UploadFile

UPLOAD_DIR = Path(os.getenv("UPLOAD_DIR", "./data/uploads"))
MAX_FILE_BYTES = 100 * 1024 * 1024  # 100MB
SPARSE_TEXT_THRESHOLD = 100  # if PDF extracts < this many chars, treat as scanned
MAX_PDF_PAGES_AS_IMAGES = 50
PDF_RENDER_DPI = 130  # JPEG @ 130 DPI balances small-text legibility vs OpenAI vision payload size
PDF_JPEG_QUALITY = 85  # good OCR quality, ~5-10x smaller than PNG
MAX_EXTRACT_CHARS = 500_000

IMAGE_TYPES = {"image/jpeg", "image/jpg", "image/png", "image/webp", "image/gif"}
PDF_TYPES = {"application/pdf"}
DOCX_TYPES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
}
XLSX_TYPES = {
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
}
# Legacy Excel 97-2003 (.xls). Browsers send "application/vnd.ms-excel" but
# often an empty/odd MIME too, so we also allow by .xls extension below.
XLS_TYPES = {"application/vnd.ms-excel"}
XLS_EXTENSIONS = {".xls"}
PPTX_TYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

# Text/code files come with many different MIME types from browsers (text/plain,
# application/octet-stream, or empty). Allow them by extension instead of MIME.
TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".log",
    ".ini", ".env", ".cfg", ".conf", ".toml",
    ".py", ".js", ".mjs", ".ts", ".tsx", ".jsx",
    ".html", ".htm", ".css", ".scss", ".sass",
    ".sql", ".yaml", ".yml", ".xml",
    ".sh", ".bash", ".zsh", ".ps1", ".bat", ".cmd",
    ".java", ".c", ".h", ".cpp", ".hpp", ".cs", ".go", ".rs",
    ".rb", ".php", ".kt", ".swift", ".dart", ".lua", ".r", ".scala",
    ".vue", ".svelte", ".gradle", ".pl", ".pm",
}

MIME_ALLOWED = IMAGE_TYPES | PDF_TYPES | DOCX_TYPES | XLSX_TYPES | XLS_TYPES | PPTX_TYPES


def _upload_root() -> Path:
    """Return the absolute upload directory path."""
    return UPLOAD_DIR.resolve()


def _ensure_dir() -> None:
    _upload_root().mkdir(parents=True, exist_ok=True)


def validate_upload_path(file_path: str | Path) -> Path:
    """
    Validate that file_path points to a real file inside UPLOAD_DIR.

    This prevents path traversal attacks such as:
    - ../../.env
    - ..\\..\\data\\chatbot.db
    - C:\\Users\\...\\secret.txt
    """
    root = _upload_root()
    requested = Path(file_path)

    # If caller passes only a filename, treat it as a file under UPLOAD_DIR.
    # Example: "abc.pdf" -> "<UPLOAD_DIR>/abc.pdf"
    if not requested.is_absolute() and requested.parent == Path("."):
        requested = root / requested

    requested = requested.resolve()

    try:
        requested.relative_to(root)
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid attachment path")

    if not requested.exists() or not requested.is_file():
        raise HTTPException(status_code=404, detail="Attachment not found")

    return requested


def _is_allowed(content_type: str, filename: str) -> bool:
    if content_type in MIME_ALLOWED:
        return True
    ext = Path(filename or "").suffix.lower()
    return ext in TEXT_EXTENSIONS or ext in XLS_EXTENSIONS


def _default_ext_for(content_type: str) -> str:
    if content_type in PDF_TYPES:
        return ".pdf"
    if content_type in DOCX_TYPES:
        return ".docx"
    if content_type in XLSX_TYPES:
        return ".xlsx"
    if content_type in XLS_TYPES:
        return ".xls"
    if content_type in PPTX_TYPES:
        return ".pptx"
    if content_type in IMAGE_TYPES:
        return ".png"
    return ".bin"


async def save_upload(upload: UploadFile) -> tuple[str, str, int, str]:
    if not _is_allowed(upload.content_type or "", upload.filename or ""):
        raise HTTPException(
            400,
            f"ไฟล์ประเภทนี้ไม่รองรับ ({upload.content_type or Path(upload.filename or '').suffix or 'unknown'})",
        )

    _ensure_dir()

    ext = Path(upload.filename or "").suffix.lower()
    if not ext:
        ext = _default_ext_for(upload.content_type or "")

    stored_name = f"{uuid.uuid4().hex}{ext}"
    file_path = (_upload_root() / stored_name).resolve()

    # Extra safety check: ensure generated path is still inside UPLOAD_DIR.
    try:
        file_path.relative_to(_upload_root())
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid upload path")

    size = 0
    try:
        with file_path.open("wb") as f:
            while chunk := await upload.read(64 * 1024):
                size += len(chunk)
                if size > MAX_FILE_BYTES:
                    f.close()
                    file_path.unlink(missing_ok=True)
                    raise HTTPException(400, "ไฟล์ใหญ่เกิน 100MB")
                f.write(chunk)
    except HTTPException:
        raise
    except Exception as e:
        file_path.unlink(missing_ok=True)
        raise HTTPException(status_code=500, detail=f"บันทึกไฟล์ไม่สำเร็จ: {e}")

    return upload.filename or stored_name, upload.content_type or "", size, str(file_path)


def is_image(content_type: str) -> bool:
    return content_type in IMAGE_TYPES


def is_pdf(content_type: str) -> bool:
    return content_type in PDF_TYPES


def is_docx(content_type: str, filename: str = "") -> bool:
    return content_type in DOCX_TYPES or filename.lower().endswith(".docx")


def is_xlsx(content_type: str, filename: str = "") -> bool:
    return content_type in XLSX_TYPES or filename.lower().endswith(".xlsx")


def is_xls(content_type: str, filename: str = "") -> bool:
    # Guard the .xlsx check first (it also ends with ...xls? no — .xlsx ≠ .xls).
    name = filename.lower()
    return (content_type in XLS_TYPES or name.endswith(".xls")) and not name.endswith(".xlsx")


def is_pptx(content_type: str, filename: str = "") -> bool:
    return content_type in PPTX_TYPES or filename.lower().endswith(".pptx")


def is_text_file(filename: str) -> bool:
    return Path(filename or "").suffix.lower() in TEXT_EXTENSIONS


def encode_image_data_url(file_path: str, content_type: str) -> str:
    safe_path = validate_upload_path(file_path)
    with safe_path.open("rb") as f:
        b64 = base64.b64encode(f.read()).decode("ascii")
    return f"data:{content_type};base64,{b64}"


def _truncate(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    return text[:max_chars] + "\n[…ตัดเนื้อหาที่เหลือเพราะยาวเกิน…]"


def extract_pdf_text(file_path: str, max_chars: int = MAX_EXTRACT_CHARS) -> str:
    try:
        safe_path = validate_upload_path(file_path)
        doc = fitz.open(str(safe_path))
        parts: list[str] = []
        total = 0

        for i, page in enumerate(doc):
            text = (page.get_text() or "").strip()
            if not text:
                continue

            chunk = f"\n[หน้า {i + 1}]\n{text}\n"
            if total + len(chunk) > max_chars:
                remaining = max_chars - total
                if remaining > 0:
                    parts.append(chunk[:remaining])
                parts.append("\n[…ตัดเนื้อหาที่เหลือเพราะยาวเกิน…]")
                break

            parts.append(chunk)
            total += len(chunk)

        doc.close()
        return "".join(parts).strip()
    except HTTPException:
        raise
    except Exception as e:
        return f"(อ่าน PDF ไม่สำเร็จ: {e})"


def extract_docx_text(file_path: str, max_chars: int = MAX_EXTRACT_CHARS) -> str:
    try:
        safe_path = validate_upload_path(file_path)

        from docx import Document

        doc = Document(str(safe_path))
        parts: list[str] = []
        total = 0

        def add(s: str) -> bool:
            nonlocal total

            if total + len(s) > max_chars:
                remaining = max_chars - total
                if remaining > 0:
                    parts.append(s[:remaining])
                parts.append("\n[…ตัดเนื้อหา…]")
                total = max_chars
                return False

            parts.append(s)
            total += len(s)
            return True

        for para in doc.paragraphs:
            t = para.text.strip()
            if not t:
                continue
            if not add(t + "\n"):
                return "".join(parts).strip()

        for t_idx, table in enumerate(doc.tables):
            if not add(f"\n[ตาราง {t_idx + 1}]\n"):
                return "".join(parts).strip()

            for row in table.rows:
                row_text = " | ".join(c.text.strip() for c in row.cells) + "\n"
                if not add(row_text):
                    return "".join(parts).strip()

        return "".join(parts).strip()
    except HTTPException:
        raise
    except Exception as e:
        return f"(อ่าน Word ไม่สำเร็จ: {e})"


def extract_xlsx_text(file_path: str, max_chars: int = MAX_EXTRACT_CHARS) -> str:
    try:
        safe_path = validate_upload_path(file_path)

        from openpyxl import load_workbook

        wb = load_workbook(str(safe_path), read_only=True, data_only=True)
        parts: list[str] = []
        total = 0

        for sheet_name in wb.sheetnames:
            if total >= max_chars:
                break

            header = f"\n=== Sheet: {sheet_name} ===\n"
            if total + len(header) > max_chars:
                parts.append("[…ตัดเนื้อหา…]")
                break

            parts.append(header)
            total += len(header)

            ws = wb[sheet_name]
            for row in ws.iter_rows(values_only=True):
                if all(v is None for v in row):
                    continue

                row_text = " | ".join("" if v is None else str(v) for v in row) + "\n"
                if total + len(row_text) > max_chars:
                    parts.append("[…ตัดเนื้อหา…]")
                    wb.close()
                    return "".join(parts).strip()

                parts.append(row_text)
                total += len(row_text)

        wb.close()
        return "".join(parts).strip()
    except HTTPException:
        raise
    except Exception as e:
        return f"(อ่าน Excel ไม่สำเร็จ: {e})"


def extract_xls_text(file_path: str, max_chars: int = MAX_EXTRACT_CHARS) -> str:
    """Legacy Excel 97-2003 (.xls) — openpyxl can't read these, use xlrd.
    Same '=== Sheet: name ===' + row layout as extract_xlsx_text."""
    try:
        safe_path = validate_upload_path(file_path)

        import xlrd

        book = xlrd.open_workbook(str(safe_path))
        parts: list[str] = []
        total = 0

        for sheet in book.sheets():
            if total >= max_chars:
                break

            header = f"\n=== Sheet: {sheet.name} ===\n"
            if total + len(header) > max_chars:
                parts.append("[…ตัดเนื้อหา…]")
                break
            parts.append(header)
            total += len(header)

            for r in range(sheet.nrows):
                values = sheet.row_values(r)
                if all(v == "" or v is None for v in values):
                    continue

                row_text = (
                    " | ".join("" if v is None else str(v) for v in values) + "\n"
                )
                if total + len(row_text) > max_chars:
                    parts.append("[…ตัดเนื้อหา…]")
                    return "".join(parts).strip()

                parts.append(row_text)
                total += len(row_text)

        return "".join(parts).strip()
    except HTTPException:
        raise
    except Exception as e:
        return f"(อ่าน Excel (.xls) ไม่สำเร็จ: {e})"


def extract_pptx_text(file_path: str, max_chars: int = MAX_EXTRACT_CHARS) -> str:
    try:
        safe_path = validate_upload_path(file_path)

        from pptx import Presentation

        prs = Presentation(str(safe_path))
        parts: list[str] = []
        total = 0

        for i, slide in enumerate(prs.slides):
            if total >= max_chars:
                break

            header = f"\n=== Slide {i + 1} ===\n"
            if total + len(header) > max_chars:
                parts.append("[…ตัดเนื้อหา…]")
                break

            parts.append(header)
            total += len(header)

            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue

                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if not t:
                        continue

                    chunk = t + "\n"
                    if total + len(chunk) > max_chars:
                        parts.append("[…ตัดเนื้อหา…]")
                        return "".join(parts).strip()

                    parts.append(chunk)
                    total += len(chunk)

        return "".join(parts).strip()
    except HTTPException:
        raise
    except Exception as e:
        return f"(อ่าน PowerPoint ไม่สำเร็จ: {e})"


def extract_text_file(file_path: str, max_chars: int = MAX_EXTRACT_CHARS) -> str:
    """Read a plain text/code file. Tries UTF-8 first, falls back to common Thai encodings."""
    safe_path = validate_upload_path(file_path)

    for enc in ("utf-8", "utf-8-sig", "cp874", "tis-620", "latin-1"):
        try:
            with safe_path.open("r", encoding=enc) as f:
                text = f.read(max_chars + 1)
            return _truncate(text, max_chars)
        except UnicodeDecodeError:
            continue
        except HTTPException:
            raise
        except Exception as e:
            return f"(อ่านไฟล์ไม่สำเร็จ: {e})"

    return "(อ่านไฟล์ไม่สำเร็จ: encoding ไม่รองรับ)"


def extract_any_text(file_path: str, content_type: str, filename: str) -> str:
    """Dispatch text extraction by file type. Returns '' for images (use vision instead)."""
    if is_pdf(content_type):
        return extract_pdf_text(file_path)
    if is_docx(content_type, filename):
        return extract_docx_text(file_path)
    if is_xlsx(content_type, filename):
        return extract_xlsx_text(file_path)
    if is_xls(content_type, filename):
        return extract_xls_text(file_path)
    if is_pptx(content_type, filename):
        return extract_pptx_text(file_path)
    if is_text_file(filename):
        return extract_text_file(file_path)
    return ""


def render_pdf_pages_as_images(
    file_path: str,
    max_pages: int = MAX_PDF_PAGES_AS_IMAGES,
    dpi: int = PDF_RENDER_DPI,
) -> list[str]:
    urls: list[str] = []

    try:
        safe_path = validate_upload_path(file_path)
        doc = fitz.open(str(safe_path))
        zoom = dpi / 72
        matrix = fitz.Matrix(zoom, zoom)

        for i, page in enumerate(doc):
            if i >= max_pages:
                break

            pix = page.get_pixmap(matrix=matrix, alpha=False)
            buf = io.BytesIO(pix.tobytes("jpeg", jpg_quality=PDF_JPEG_QUALITY))
            b64 = base64.b64encode(buf.getvalue()).decode("ascii")
            urls.append(f"data:image/jpeg;base64,{b64}")

        doc.close()
    except HTTPException:
        raise
    except Exception:
        pass

    return urls


def count_pdf_pages(file_path: str) -> int:
    """Total page count of a PDF (0 on error). Used to warn when a scanned PDF
    exceeds MAX_PDF_PAGES_AS_IMAGES and only its first pages get read."""
    try:
        safe_path = validate_upload_path(file_path)
        doc = fitz.open(str(safe_path))
        n = len(doc)
        doc.close()
        return n
    except Exception:
        return 0


def pdf_truncation_note(filename: str, file_path: str) -> str | None:
    """If a scanned PDF has more pages than the vision cap, return a Thai notice
    (with the reason) explaining only the first N pages were read. Else None."""
    total = count_pdf_pages(file_path)
    if total <= MAX_PDF_PAGES_AS_IMAGES:
        return None
    return (
        f"⚠️ หมายเหตุ: ไฟล์ \"{filename}\" เป็น PDF สแกน (รูปภาพ) มีทั้งหมด {total} หน้า — "
        f"ระบบอ่านไฟล์สแกนด้วย AI ภาพ ซึ่งจำกัดที่ {MAX_PDF_PAGES_AS_IMAGES} หน้าแรกต่อไฟล์ "
        f"จึงประมวลผลเฉพาะหน้า 1-{MAX_PDF_PAGES_AS_IMAGES} (หน้า {MAX_PDF_PAGES_AS_IMAGES + 1}-{total} ยังไม่ถูกอ่าน) "
        f"หากต้องการครบ กรุณาแบ่งไฟล์เป็นช่วงละไม่เกิน {MAX_PDF_PAGES_AS_IMAGES} หน้า"
    )


# Asking to translate/transcribe a WHOLE long document in chat produces
# incomplete or paraphrased output (the chat answer is token-capped and the
# model drifts into summary). Detect that intent so we can warn and point the
# user to the dedicated page-by-page translation tool.
_TRANSLATE_INTENT_RE = re.compile(r"แปล|ถอดความ|translate|transcrib", re.IGNORECASE)

# Below these sizes a chat translation is fine (short doc fits one answer).
_LONG_TEXT_CHARS = 6000
_LONG_IMAGE_PAGES = 3


def long_translation_warning(question: str, text_chars: int, image_pages: int) -> str | None:
    """Return a Thai warning if the user asks to translate a LONG attached
    document in chat (where the answer would be incomplete/paraphrased). Else None."""
    if not question or not _TRANSLATE_INTENT_RE.search(question):
        return None
    if text_chars <= _LONG_TEXT_CHARS and image_pages < _LONG_IMAGE_PAGES:
        return None
    return (
        "⚠️ หมายเหตุ: การแปล/ถอดความเอกสารยาวผ่านแชทอาจ **ไม่ครบหรือไม่ตรงต้นฉบับ 100%** "
        "เพราะระบบแชทตอบได้จำกัดต่อรอบ และอาจสรุป/เรียบเรียงใหม่โดยไม่ได้ตั้งใจ "
        "หากต้องการคำแปลที่ครบทุกหน้า ตรงต้นฉบับ พร้อมไฟล์ Word/PDF "
        "แนะนำให้ใช้เครื่องมือแปลเอกสารเฉพาะทาง (แจ้งผู้ดูแลระบบได้เลย)"
    )


def is_text_sparse(extracted: str) -> bool:
    return len(extracted.strip()) < SPARSE_TEXT_THRESHOLD