"""
Bulk-import knowledge base content from department files.

Supports two input modes:

1. Q&A JSON (cleanest, preferred for FAQ-style content from HR/Sales/etc.):

       python import_doc.py --qa-json hr_faq.json --dept HR

   The JSON must be a list of {"question": "...", "answer": "..."} objects.

2. Document chunking (for long manuals / SOPs / PDFs / Word):

       python import_doc.py --file kb_sources/production/komori_manual.pdf --dept Production
       python import_doc.py --dir  kb_sources/HR                            --dept HR

   The script extracts text from the file (PDF/DOCX/XLSX/PPTX), splits it
   into ~400-token chunks at paragraph boundaries, and imports each chunk
   as a knowledge entry.

Optional metadata Excel (--metadata kb_metadata.xlsx):

   Recognized header columns (case-insensitive, spaces normalized):
       file_name             [required key]
       department            (overrides --dept for matched files)
       document_type         (SOP, Policy, Manual, FAQ, ...)
       confidentiality_level (public | internal | confidential)
       allowed_groups        (comma-separated)
       last_updated          (date)
       note                  (freeform)

   Each chunk imported from a file gets the file's metadata attached via
   UPDATE after add_knowledge() returns the new row id. Run
   migrate_kb_metadata.py once to add the underlying columns.

Both modes are idempotent: rows whose `question` already exists are skipped.

The script requires an admin user to exist in the chatbot DB (it stamps
`created_by`). Promote a user to admin first if you haven't:

    docker exec siriwattana-postgres-test psql -U chatbot -d chatbot_test \\
      -c "UPDATE users SET role='admin' WHERE username='<name>';"

Runs against whichever DB_ENGINE is configured in backend/.env (sqlite or postgres).
"""

import argparse
import json
import os
import re
import sys
from pathlib import Path

from dotenv import load_dotenv

load_dotenv()

from rag import add_knowledge

DB_ENGINE = os.getenv("DB_ENGINE", "sqlite").strip().lower()
USE_PG = DB_ENGINE in {"postgres", "postgresql", "pg"}

CHUNK_TARGET_CHARS = 1600  # rough proxy for ~400 tokens of Thai text
CHUNK_MIN_CHARS = 200      # don't import tiny scraps


METADATA_COLS = [
    "source_file",
    "source_dept",
    "document_type",
    "confidentiality",
    "allowed_groups",
    "last_updated",
    "note",
]


HEADER_ALIASES = {
    "file_name": "source_file",
    "filename": "source_file",
    "file": "source_file",
    "department": "source_dept",
    "dept": "source_dept",
    "document_type": "document_type",
    "doc_type": "document_type",
    "type": "document_type",
    "confidentiality_level": "confidentiality",
    "confidentiality": "confidentiality",
    "confid": "confidentiality",
    "allowed_groups": "allowed_groups",
    "groups": "allowed_groups",
    "last_updated": "last_updated",
    "updated": "last_updated",
    "note": "note",
    "notes": "note",
    "remark": "note",
}


VALID_CONFIDENTIALITY = {"public", "internal", "confidential"}


# ─────────────────────── DB helpers (reused pattern from seed.py) ───────────────────────


def fetch_admin():
    if USE_PG:
        from db_pg import connect_pg

        conn = connect_pg()
        try:
            with conn.cursor() as cur:
                cur.execute(
                    "SELECT id, username FROM users WHERE role = 'admin' "
                    "ORDER BY id LIMIT 1"
                )
                row = cur.fetchone()
                return {"id": row[0], "username": row[1]} if row else None
        finally:
            conn.close()

    from db import get_db

    conn = get_db()
    row = conn.execute(
        "SELECT id, username FROM users WHERE role = 'admin' ORDER BY id LIMIT 1"
    ).fetchone()
    return dict(row) if row else None


def question_exists(question: str) -> bool:
    if USE_PG:
        from db_pg import connect_pg

        conn = connect_pg()
        try:
            with conn.cursor() as cur:
                cur.execute(
                    "SELECT id FROM knowledge WHERE question = %s", (question,)
                )
                return cur.fetchone() is not None
        finally:
            conn.close()

    from db import get_db

    conn = get_db()
    return (
        conn.execute(
            "SELECT id FROM knowledge WHERE question = ?", (question,)
        ).fetchone()
        is not None
    )


def update_metadata(knowledge_id: int, meta: dict) -> None:
    """Apply metadata to an existing knowledge row.

    Silently ignores columns that don't exist in the schema yet (so this
    script keeps working even before migrate_kb_metadata.py has been run).
    """
    if not meta:
        return

    set_pairs = [(col, meta[col]) for col in METADATA_COLS if col in meta and meta[col] is not None]
    if not set_pairs:
        return

    if USE_PG:
        from db_pg import connect_pg

        placeholders = ", ".join(f"{col} = %s" for col, _ in set_pairs)
        values = [v for _, v in set_pairs] + [knowledge_id]
        conn = connect_pg()
        try:
            with conn.cursor() as cur:
                try:
                    cur.execute(
                        f"UPDATE knowledge SET {placeholders} WHERE id = %s", values
                    )
                except Exception as exc:
                    # Likely: a column doesn't exist yet. Don't crash the whole import.
                    print(f"      (metadata update skipped: {exc})")
            conn.commit()
        finally:
            conn.close()
        return

    from db import get_db

    placeholders = ", ".join(f"{col} = ?" for col, _ in set_pairs)
    values = [v for _, v in set_pairs] + [knowledge_id]
    conn = get_db()
    try:
        conn.execute(f"UPDATE knowledge SET {placeholders} WHERE id = ?", values)
        conn.commit()
    except Exception as exc:
        print(f"      (metadata update skipped: {exc})")


def load_metadata_excel(xlsx_path: Path) -> dict[str, dict]:
    """Read an Excel metadata sheet and return {file_name_lower: meta_dict}.

    Header matching is case-insensitive; spaces and dashes are normalized to
    underscores. Recognized aliases are listed in HEADER_ALIASES.
    """
    from openpyxl import load_workbook

    wb = load_workbook(str(xlsx_path), data_only=True, read_only=True)
    ws = wb.active
    if ws is None:
        return {}

    rows = ws.iter_rows(values_only=True)
    try:
        header_row = next(rows)
    except StopIteration:
        return {}

    def normalize(h):
        if h is None:
            return None
        return str(h).strip().lower().replace(" ", "_").replace("-", "_")

    headers = [normalize(h) for h in header_row]
    col_map: dict[int, str] = {}
    for idx, h in enumerate(headers):
        if not h:
            continue
        canonical = HEADER_ALIASES.get(h)
        if canonical:
            col_map[idx] = canonical

    if "source_file" not in col_map.values():
        print(
            "  WARNING: metadata sheet has no 'file_name' column — every row will be ignored."
        )
        return {}

    out: dict[str, dict] = {}
    for row in rows:
        meta: dict = {}
        for idx, canonical in col_map.items():
            if idx >= len(row):
                continue
            val = row[idx]
            if val is None or (isinstance(val, str) and not val.strip()):
                continue
            if canonical == "confidentiality":
                v = str(val).strip().lower()
                if v.endswith("_level"):
                    v = v.replace("_level", "")
                if v not in VALID_CONFIDENTIALITY:
                    print(f"    WARNING: confidentiality '{val}' not in {VALID_CONFIDENTIALITY}, skipping")
                    continue
                meta[canonical] = v
            elif canonical == "last_updated":
                # openpyxl gives us a datetime/date directly; ISO-format strings are fine too.
                meta[canonical] = str(val)[:10] if val else None
            else:
                meta[canonical] = str(val).strip() if not isinstance(val, str) else val.strip()

        fname = meta.get("source_file")
        if not fname:
            continue
        out[fname.lower()] = meta

    print(f"  Loaded metadata for {len(out)} file(s) from {xlsx_path.name}")
    return out


def metadata_for_file(
    path: Path,
    fallback_dept: str | None,
    metadata_map: dict[str, dict],
) -> dict:
    """Look up metadata for `path` in metadata_map, falling back to defaults."""
    base = metadata_map.get(path.name.lower(), {}).copy()
    base.setdefault("source_file", path.name)
    if fallback_dept and not base.get("source_dept"):
        base["source_dept"] = fallback_dept
    base.setdefault("confidentiality", "internal")
    return base


# ─────────────────────── File extraction ───────────────────────


def extract_text_pdf(path: Path) -> str:
    import pypdf

    reader = pypdf.PdfReader(str(path))
    parts = []
    for page in reader.pages:
        try:
            parts.append(page.extract_text() or "")
        except Exception:
            continue
    return "\n\n".join(parts)


def extract_text_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def extract_text_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), data_only=True, read_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"=== Sheet: {sheet_name} ===")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_text_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            parts.append(f"=== Slide {i} ===\n" + "\n".join(slide_text))
    return "\n\n".join(parts)


EXTRACTORS = {
    ".pdf": extract_text_pdf,
    ".docx": extract_text_docx,
    ".xlsx": extract_text_xlsx,
    ".pptx": extract_text_pptx,
}


def extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext not in EXTRACTORS:
        raise ValueError(
            f"Unsupported file type: {ext} (supported: {', '.join(EXTRACTORS)})"
        )
    return EXTRACTORS[ext](path)


# ─────────────────────── Chunking ───────────────────────


def chunk_text(text: str, target: int = CHUNK_TARGET_CHARS) -> list[str]:
    """Split text into ~target-char chunks at paragraph boundaries.

    Falls back to splitting on single newlines, then sentence boundaries,
    if paragraphs alone are too long.
    """
    text = re.sub(r"\n{3,}", "\n\n", text.strip())
    if not text:
        return []

    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]

    chunks: list[str] = []
    buffer = ""

    def flush():
        nonlocal buffer
        if buffer.strip():
            chunks.append(buffer.strip())
        buffer = ""

    for para in paragraphs:
        if len(para) > target:
            flush()
            # Para itself is huge — split on newlines, then sentences.
            for line in para.split("\n"):
                line = line.strip()
                if not line:
                    continue
                if len(line) > target:
                    sentences = re.split(r"(?<=[.!?。．！？ฯ])\s+", line)
                    for s in sentences:
                        s = s.strip()
                        if not s:
                            continue
                        if len(buffer) + len(s) + 1 > target:
                            flush()
                        buffer = (buffer + " " + s).strip() if buffer else s
                else:
                    if len(buffer) + len(line) + 1 > target:
                        flush()
                    buffer = (buffer + "\n" + line).strip() if buffer else line
            flush()
            continue

        if len(buffer) + len(para) + 2 > target:
            flush()
        buffer = (buffer + "\n\n" + para).strip() if buffer else para

    flush()
    return [c for c in chunks if len(c) >= CHUNK_MIN_CHARS]


# ─────────────────────── Import paths ───────────────────────


def import_qa_pairs(
    pairs: list[dict],
    admin: dict,
    dept: str | None,
    source_label: str,
    base_meta: dict | None = None,
) -> tuple[int, int]:
    added = 0
    skipped = 0
    for i, item in enumerate(pairs, 1):
        question = (item.get("question") or "").strip()
        answer = (item.get("answer") or "").strip()
        if not question or not answer:
            print(f"  [{i:>3}/{len(pairs)}] SKIP (empty Q or A)")
            skipped += 1
            continue
        if question_exists(question):
            print(f"  [{i:>3}/{len(pairs)}] SKIP (exists): {question[:60]}")
            skipped += 1
            continue
        kid = add_knowledge(question, answer, admin["id"], source=source_label)
        added += 1
        print(f"  [{i:>3}/{len(pairs)}] ADD  id={kid}: {question[:60]}")

        # Per-item metadata can override the base meta (e.g. confidentiality on a single Q).
        meta = dict(base_meta or {})
        for key in METADATA_COLS:
            if key in item and item[key] is not None:
                meta[key] = item[key]
        update_metadata(kid, meta)

    return added, skipped


def import_doc_chunks(
    path: Path,
    admin: dict,
    dept: str | None,
    source_label: str,
    metadata_map: dict[str, dict] | None = None,
) -> tuple[int, int]:
    print(f"  Extracting text from {path.name} ...")
    text = extract_text(path)
    if not text.strip():
        print(f"  (empty text — nothing to import)")
        return 0, 0

    chunks = chunk_text(text)
    print(f"  Got {len(chunks)} chunks (after filtering < {CHUNK_MIN_CHARS} chars)")

    meta_for_chunks = metadata_for_file(path, dept, metadata_map or {})
    eff_dept = meta_for_chunks.get("source_dept") or dept

    added = 0
    skipped = 0
    file_tag = f"[ที่มา: {path.name}{f' / {eff_dept}' if eff_dept else ''}]"

    for i, chunk in enumerate(chunks, 1):
        # Use the first non-empty line as a heuristic title; if too long, truncate.
        first_line = next((ln.strip() for ln in chunk.splitlines() if ln.strip()), chunk)
        title = first_line[:120]
        # The "question" field carries the title + file tag so admin can spot the source later.
        question = f"{title} {file_tag}".strip()
        answer = f"{chunk}\n\n{file_tag}"

        if question_exists(question):
            print(f"  [{i:>3}/{len(chunks)}] SKIP (exists): {title[:60]}")
            skipped += 1
            continue

        kid = add_knowledge(question, answer, admin["id"], source=source_label)
        added += 1
        print(f"  [{i:>3}/{len(chunks)}] ADD  id={kid}: {title[:60]}")
        update_metadata(kid, meta_for_chunks)

    return added, skipped


# ─────────────────────── CLI ───────────────────────


def main() -> None:
    parser = argparse.ArgumentParser(description="Import knowledge base content into the chatbot DB.")
    g = parser.add_mutually_exclusive_group(required=True)
    g.add_argument("--qa-json", type=Path, help="JSON file with [{question, answer}, ...]")
    g.add_argument("--file", type=Path, help="Single PDF/DOCX/XLSX/PPTX to chunk-import")
    g.add_argument("--dir", type=Path, help="Directory — import every supported file inside (non-recursive)")
    parser.add_argument("--dept", help="Department label, recorded in the `source` column (e.g. HR, Production)")
    parser.add_argument("--metadata", type=Path, help="Optional Excel sheet with per-file metadata (see header docs above)")
    parser.add_argument("--dry-run", action="store_true", help="Show what would be imported, don't write to DB")
    args = parser.parse_args()

    admin = fetch_admin()
    if not admin:
        print("ERROR: ยังไม่มี admin user — promote คนหนึ่งเป็น admin ก่อน")
        sys.exit(1)

    source_label = f"kb_import:{args.dept}" if args.dept else "kb_import"
    print(f"DB_ENGINE:     {DB_ENGINE}")
    print(f"Admin:         {admin['username']} (id={admin['id']})")
    print(f"Dept:          {args.dept or '(none)'}")
    print(f"Source label:  {source_label}")
    print(f"Dry-run:       {args.dry_run}")
    print()

    if args.dry_run:
        # Make BOTH add_knowledge and update_metadata into no-ops for dry-run.
        # Important: `from rag import add_knowledge` binds the name in THIS module's
        # globals, so patching only rag.add_knowledge does nothing — we have to
        # overwrite the local binding too. Also patch rag's copy for completeness.
        import rag

        noop_add = lambda q, a, uid, source="admin": 0  # type: ignore
        rag.add_knowledge = noop_add
        globals()["add_knowledge"] = noop_add
        globals()["update_metadata"] = lambda kid, meta: None

    metadata_map: dict[str, dict] = {}
    if args.metadata:
        if not args.metadata.exists():
            print(f"ERROR: metadata file not found: {args.metadata}")
            sys.exit(1)
        try:
            metadata_map = load_metadata_excel(args.metadata)
        except Exception as exc:
            print(f"ERROR reading metadata Excel: {exc}")
            sys.exit(1)

    total_added = 0
    total_skipped = 0

    if args.qa_json:
        if not args.qa_json.exists():
            print(f"ERROR: file not found: {args.qa_json}")
            sys.exit(1)
        pairs = json.loads(args.qa_json.read_text(encoding="utf-8"))
        if not isinstance(pairs, list):
            print("ERROR: Q&A JSON must be a list of {question, answer} objects")
            sys.exit(1)
        # Q&A JSON inherits metadata from the JSON file's name if the Excel has a row for it.
        base_meta = metadata_for_file(args.qa_json, args.dept, metadata_map) if metadata_map else {
            "source_file": args.qa_json.name,
            "source_dept": args.dept,
            "confidentiality": "internal",
        }
        print(f"Importing {len(pairs)} Q&A pair(s) from {args.qa_json.name}")
        added, skipped = import_qa_pairs(pairs, admin, args.dept, source_label, base_meta=base_meta)
        total_added += added
        total_skipped += skipped

    elif args.file:
        if not args.file.exists():
            print(f"ERROR: file not found: {args.file}")
            sys.exit(1)
        print(f"Importing chunks from {args.file}")
        added, skipped = import_doc_chunks(args.file, admin, args.dept, source_label, metadata_map)
        total_added += added
        total_skipped += skipped

    elif args.dir:
        if not args.dir.is_dir():
            print(f"ERROR: not a directory: {args.dir}")
            sys.exit(1)
        # Recursive: pick up files inside subfolders too.
        files = sorted(
            p for p in args.dir.rglob("*")
            if p.is_file() and p.suffix.lower() in EXTRACTORS
        )
        print(f"Found {len(files)} supported file(s) under {args.dir} (recursive)")
        for path in files:
            # Show relative path so subfolder context is visible.
            try:
                rel = path.relative_to(args.dir)
            except ValueError:
                rel = path
            print(f"\n--- {rel} ---")
            try:
                added, skipped = import_doc_chunks(path, admin, args.dept, source_label, metadata_map)
                total_added += added
                total_skipped += skipped
            except Exception as exc:
                print(f"  ERROR processing {path.name}: {exc}")

    print()
    print(f"=== สรุป ===")
    print(f"เพิ่มใหม่:  {total_added}")
    print(f"ข้าม:       {total_skipped}")
    if args.dry_run:
        print("(dry-run — ไม่มีข้อมูลถูกเขียนลง DB)")


if __name__ == "__main__":
    main()
